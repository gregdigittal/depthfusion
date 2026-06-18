"""PPTX document parser (T-595).

Parses Microsoft PowerPoint .pptx files using python-pptx.  Each slide
produces one DocumentRecord.  Speaker notes are appended to the slide
content when present.

Lazy import: python-pptx is optional.  If it is not installed this parser
returns an empty list for every call rather than raising.
"""
from __future__ import annotations

import io
from datetime import datetime, timezone

from depthfusion.parsers.documents.base import DocumentParser, DocumentRecord  # noqa: F401

try:
    from pptx import Presentation
    from pptx.util import Pt  # noqa: F401 — imported to confirm pptx is usable

    _pptx_available = True
except ImportError:  # pragma: no cover
    _pptx_available = False

_MIME_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


class PptxParser:
    """DocumentParser for PowerPoint .pptx files."""

    name: str = "pptx"
    supported_mime_types: list[str] = [_MIME_PPTX]

    def parse(self, source_id: str, data: bytes) -> list[DocumentRecord]:
        """Parse *data* as a PowerPoint presentation.

        Returns one :class:`DocumentRecord` per slide, or an empty list
        if python-pptx is unavailable, *data* is empty, or any error occurs.
        """
        if not _pptx_available or not data:
            return []

        try:
            return self._parse_presentation(source_id, data)
        except Exception:  # noqa: BLE001
            return []

    def _parse_presentation(self, source_id: str, data: bytes) -> list[DocumentRecord]:
        prs = Presentation(io.BytesIO(data))
        parse_ts = datetime.now(timezone.utc).isoformat().replace("+00:00", "Z")
        records: list[DocumentRecord] = []

        for i, slide in enumerate(prs.slides):
            # ── Title ──────────────────────────────────────────────────────
            title_text = ""
            if slide.shapes.title is not None:
                title_text = (slide.shapes.title.text or "").strip()
            if not title_text:
                title_text = f"Slide {i + 1}"

            # ── Body text chunks (one per text-bearing shape) ──────────────
            chunks: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if text:
                    chunks.append(text)

            # ── Speaker notes ──────────────────────────────────────────────
            notes_text = ""
            try:
                notes_slide = slide.notes_slide
                if notes_slide is not None:
                    raw = notes_slide.notes_text_frame.text.strip()
                    if raw:
                        notes_text = raw
            except Exception:  # noqa: BLE001
                pass

            # ── Assemble content ───────────────────────────────────────────
            content_parts = list(chunks)
            if notes_text:
                content_parts.append(f"Speaker notes: {notes_text}")

            content = "\n".join(content_parts)

            records.append(
                DocumentRecord(
                    source_id=source_id,
                    title=title_text,
                    content=content,
                    chunks=chunks,
                    heading_path=[title_text],
                    mime_type=_MIME_PPTX,
                    parse_timestamp=parse_ts,
                )
            )

        return records


__all__ = ["PptxParser"]
