"""Golden corpus + snapshot tests for xlsx and pptx document parsers (T-596).

Tests use synthetic in-memory files built with the respective libraries so
no real fixture files are needed.  All test classes are self-contained.
"""
from __future__ import annotations

import io
from datetime import datetime

import docx as python_docx
import openpyxl
import pytest
from pptx import Presentation
from pptx.util import Inches

from depthfusion.parsers.documents import (
    DocxParser,
    OcrParser,
    PdfParser,
    PptxParser,
    XlsxParser,
    get_registry,
)
from depthfusion.parsers.documents.docx import _MIME_DOCX
from depthfusion.parsers.documents.ocr import _MIME_PNG
from depthfusion.parsers.documents.pdf import _MIME_PDF
from depthfusion.parsers.documents.pptx import _MIME_PPTX
from depthfusion.parsers.documents.xlsx import _MIME_XLSX

# ──────────────────────────────────────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────────────────────────────────────

def _make_xlsx(*sheet_specs: list[tuple]) -> bytes:
    """Build an in-memory .xlsx workbook.

    Each positional arg is a list of row-tuples for one sheet::

        _make_xlsx(
            [("Name", "Value"), ("Alice", 42)],   # Sheet1
            [("X", "Y")],                          # Sheet2
        )
    """
    wb = openpyxl.Workbook()
    wb.remove(wb.active)  # remove default empty sheet

    for idx, rows in enumerate(sheet_specs):
        ws = wb.create_sheet(title=f"Sheet{idx + 1}")
        for row in rows:
            ws.append(list(row))

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx(*slides: dict) -> bytes:
    """Build an in-memory .pptx presentation.

    Each positional arg is a dict with optional keys:
        title   (str): slide title text
        body    (str): body placeholder text
        notes   (str): speaker-notes text
    """
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # Blank

    for spec in slides:
        title_text = spec.get("title", "")
        body_text = spec.get("body", "")
        notes_text = spec.get("notes", "")

        if title_text and body_text:
            layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title_text
            slide.placeholders[1].text = body_text
        elif title_text:
            layout = prs.slide_layouts[0]  # Title Slide
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title_text
        else:
            slide = prs.slides.add_slide(blank_layout)
            # Add a text box for body content if provided
            if body_text:
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
                txBox.text_frame.text = body_text

        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_docx(*sections: dict) -> bytes:
    doc = python_docx.Document()
    for spec in sections:
        if spec.get("heading"):
            doc.add_heading(spec["heading"], level=1)
        if spec.get("subheading"):
            doc.add_heading(spec["subheading"], level=2)
        if spec.get("body"):
            doc.add_paragraph(spec["body"])
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


MINIMAL_PDF = (
    b"%PDF-1.4\n"
    b"1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj\n"
    b"2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj\n"
    b"3 0 obj << /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]\n"
    b"  /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >> endobj\n"
    b"4 0 obj << /Length 44 >>\nstream\n"
    b"BT /F1 12 Tf 100 700 Td (Hello PDF) Tj ET\n"
    b"endstream\nendobj\n"
    b"5 0 obj << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> endobj\n"
    b"xref\n0 6\n"
    b"0000000000 65535 f \n"
    b"0000000009 00000 n \n"
    b"0000000058 00000 n \n"
    b"0000000115 00000 n \n"
    b"0000000266 00000 n \n"
    b"0000000360 00000 n \n"
    b"trailer << /Size 6 /Root 1 0 R >>\n"
    b"startxref\n441\n%%EOF\n"
)


def _make_pdf(num_pages: int = 1) -> bytes:
    try:
        from fpdf import FPDF

        pdf = FPDF()
        for i in range(num_pages):
            pdf.add_page()
            pdf.set_font("Helvetica", size=12)
            pdf.cell(0, 10, f"Hello PDF page {i + 1}")
        output = pdf.output()
        return output if isinstance(output, bytes) else bytes(output)
    except ImportError:
        pass

    try:
        from reportlab.pdfgen import canvas

        buf = io.BytesIO()
        c = canvas.Canvas(buf)
        for i in range(num_pages):
            c.drawString(100, 700, f"Hello PDF page {i + 1}")
            if i < num_pages - 1:
                c.showPage()
        c.save()
        return buf.getvalue()
    except ImportError:
        pass

    return MINIMAL_PDF


# ──────────────────────────────────────────────────────────────────────────────
# XlsxParser tests
# ──────────────────────────────────────────────────────────────────────────────

class TestXlsxParser:
    def setup_method(self) -> None:
        self.parser = XlsxParser()

    def test_empty_bytes_returns_empty_list(self) -> None:
        assert self.parser.parse("doc-empty", b"") == []

    def test_single_sheet_parsed(self) -> None:
        data = _make_xlsx([("Name", "Value", "Type"), ("Alice", 42, "admin"), ("Bob", 7, "user")])
        records = self.parser.parse("doc-1", data)

        assert len(records) == 1
        rec = records[0]
        assert rec.source_id == "doc-1"
        assert "Alice" in rec.content
        assert "42" in rec.content
        assert rec.mime_type == _MIME_XLSX
        assert rec.parse_timestamp.endswith("Z")

    def test_headers_inferred(self) -> None:
        data = _make_xlsx([("Name", "Value", "Type"), ("Alice", 1, "a")])
        records = self.parser.parse("doc-headers", data)

        assert len(records) == 1
        rec = records[0]
        # First row all non-empty → header row text should appear in content
        assert "Name | Value | Type" in rec.content

    def test_multi_sheet(self) -> None:
        data = _make_xlsx(
            [("A", "B"), (1, 2)],
            [("X", "Y", "Z"), (10, 20, 30)],
        )
        records = self.parser.parse("doc-multi", data)

        assert len(records) == 2
        titles = {r.title for r in records}
        assert "Sheet: Sheet1" in titles
        assert "Sheet: Sheet2" in titles

    def test_error_returns_empty_list(self) -> None:
        records = self.parser.parse("doc-bad", b"not an xlsx file")
        assert records == []

    def test_chunks_non_empty(self) -> None:
        rows = [("H1", "H2")] + [(str(i), str(i * 2)) for i in range(10)]
        data = _make_xlsx(rows)
        records = self.parser.parse("doc-chunks", data)

        assert len(records) == 1
        assert len(records[0].chunks) >= 1

    def test_heading_path_contains_sheet_title(self) -> None:
        data = _make_xlsx([("Col",), ("val",)])
        records = self.parser.parse("doc-hp", data)

        assert records[0].heading_path[0] == "Sheet1"

    def test_title_format(self) -> None:
        data = _make_xlsx([("A",), ("1",)])
        records = self.parser.parse("doc-title", data)

        assert records[0].title == "Sheet: Sheet1"


# ──────────────────────────────────────────────────────────────────────────────
# PptxParser tests
# ──────────────────────────────────────────────────────────────────────────────

class TestPptxParser:
    def setup_method(self) -> None:
        self.parser = PptxParser()

    def test_empty_bytes_returns_empty_list(self) -> None:
        assert self.parser.parse("pptx-empty", b"") == []

    def test_single_slide_title(self) -> None:
        data = _make_pptx({"title": "My Slide"})
        records = self.parser.parse("pptx-1", data)

        assert len(records) == 1
        assert records[0].title == "My Slide"

    def test_speaker_notes_included(self) -> None:
        data = _make_pptx({"title": "Notes Slide", "notes": "See me"})
        records = self.parser.parse("pptx-notes", data)

        assert len(records) == 1
        assert "See me" in records[0].content

    def test_error_returns_empty_list(self) -> None:
        records = self.parser.parse("pptx-bad", b"not a pptx file")
        assert records == []

    def test_multiple_slides(self) -> None:
        data = _make_pptx(
            {"title": "Slide One"},
            {"title": "Slide Two"},
        )
        records = self.parser.parse("pptx-multi", data)

        assert len(records) == 2
        titles = [r.title for r in records]
        assert "Slide One" in titles
        assert "Slide Two" in titles

    def test_fallback_title_when_no_title_shape(self) -> None:
        # A blank layout slide has no title shape, so fallback "Slide N" applies.
        data = _make_pptx({"body": "just body, no title"})
        records = self.parser.parse("pptx-notitle", data)

        assert len(records) == 1
        assert records[0].title == "Slide 1"

    def test_mime_type(self) -> None:
        data = _make_pptx({"title": "T"})
        records = self.parser.parse("pptx-mime", data)

        assert records[0].mime_type == _MIME_PPTX

    def test_parse_timestamp_iso_format(self) -> None:
        data = _make_pptx({"title": "T"})
        records = self.parser.parse("pptx-ts", data)

        ts = records[0].parse_timestamp
        assert ts.endswith("Z")
        assert "T" in ts  # ISO 8601 separator

    def test_heading_path(self) -> None:
        data = _make_pptx({"title": "My Slide"})
        records = self.parser.parse("pptx-hp", data)

        assert records[0].heading_path == ["My Slide"]


class TestDocxParser:
    def test_empty_bytes_returns_empty_list(self) -> None:
        assert DocxParser().parse("x.docx", b"") == []

    def test_single_heading_and_body(self) -> None:
        data = _make_docx({"heading": "Intro", "body": "Some body text"})
        records = DocxParser().parse("test.docx", data)
        assert len(records) >= 1
        r = records[0]
        assert r.heading_path
        assert r.content and len(r.content) > 0

    def test_heading_path_hierarchy(self) -> None:
        data = _make_docx(
            {"heading": "Chapter 1", "subheading": "Section 1.1", "body": "Details here"}
        )
        records = DocxParser().parse("test.docx", data)
        assert len(records) >= 1
        # At least one record should have the hierarchy separator in its heading_path
        paths = [path for r in records for path in r.heading_path if path]
        assert any("▸" in path for path in paths)

    def test_multiple_headings_produce_multiple_records(self) -> None:
        data = _make_docx(
            {"heading": "Alpha", "body": "First"},
            {"heading": "Beta", "body": "Second"},
            {"heading": "Gamma", "body": "Third"},
        )
        records = DocxParser().parse("test.docx", data)
        assert len(records) >= 3

    def test_error_returns_empty_list(self) -> None:
        assert DocxParser().parse("x.docx", b"not-a-docx-file") == []

    def test_mime_type(self) -> None:
        mime_types = DocxParser().supported_mime_types
        assert _MIME_DOCX in mime_types

    def test_parse_timestamp_iso_format(self) -> None:
        import re

        data = _make_docx({"heading": "TS Test", "body": "body"})
        records = DocxParser().parse("test.docx", data)
        assert len(records) >= 1
        ts = records[0].parse_timestamp
        assert ts is not None
        iso_pattern = r"^\d{4}-\d{2}-\d{2}T\d{2}:\d{2}:\d{2}"
        assert re.match(iso_pattern, ts), f"timestamp {ts!r} does not match ISO 8601"

    def test_document_with_no_headings(self) -> None:
        data = _make_docx({"body": "Just a plain paragraph with no heading"})
        records = DocxParser().parse("test.docx", data)
        assert len(records) >= 1
        assert any(r.content and "plain paragraph" in r.content for r in records)


class TestPdfParser:
    def test_empty_bytes_returns_empty_list(self) -> None:
        assert PdfParser().parse("x", b"") == []

    def test_single_page_produces_record(self) -> None:
        records = PdfParser().parse("pdf-1", _make_pdf(1))

        assert len(records) >= 1

    def test_page_count_matches_records(self) -> None:
        records = PdfParser().parse("pdf-2", _make_pdf(2))
        if not records:
            pytest.skip("pdfplumber could not parse the generated test PDF")

        assert len(records) == 2

    def test_heading_path_contains_page_number(self) -> None:
        records = PdfParser().parse("pdf-heading", _make_pdf(1))
        if not records:
            pytest.skip("PDF parser dependency could not parse the generated test PDF")

        heading_text = " ".join(records[0].heading_path).lower()
        assert "page" in heading_text or "1" in heading_text

    def test_document_id_set_correctly(self) -> None:
        records = PdfParser().parse("pdf-doc-id", _make_pdf(1))
        if not records:
            pytest.skip("PDF parser dependency could not parse the generated test PDF")

        assert records[0].source_id == "pdf-doc-id"

    def test_mime_type(self) -> None:
        assert _MIME_PDF in PdfParser().supported_mime_types

    def test_parse_timestamp_iso_format(self) -> None:
        records = PdfParser().parse("pdf-ts", _make_pdf(1))
        if not records:
            pytest.skip("PDF parser dependency could not parse the generated test PDF")

        ts = records[0].parse_timestamp
        assert ts.endswith("Z")
        datetime.fromisoformat(ts.replace("Z", "+00:00"))

    def test_error_returns_empty_list(self) -> None:
        assert PdfParser().parse("pdf-bad", b"not a pdf") == []


class TestOcrParser:
    def test_disabled_ocr_returns_empty_list(self, monkeypatch: pytest.MonkeyPatch) -> None:
        monkeypatch.delenv("DEPTHFUSION_OCR_ENABLED", raising=False)

        assert OcrParser().parse("x", b"") == []

    def test_mime_type(self) -> None:
        assert _MIME_PNG in OcrParser().supported_mime_types

    def test_disabled_ocr_returns_empty_for_png_bytes(
        self,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        monkeypatch.delenv("DEPTHFUSION_OCR_ENABLED", raising=False)

        assert OcrParser().parse("y", b"\x89PNG\r\n\x1a\n") == []

    def test_ocr_parser_name(self) -> None:
        assert getattr(OcrParser(), "name") == "ocr"


# ──────────────────────────────────────────────────────────────────────────────
# Registry integration tests
# ──────────────────────────────────────────────────────────────────────────────

class TestRegistryIntegration:
    def test_xlsx_registered(self) -> None:
        parser = get_registry().get(
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )
        assert parser is not None
        assert isinstance(parser, XlsxParser)

    def test_xlsx_xls_alias_registered(self) -> None:
        parser = get_registry().get("application/vnd.ms-excel")
        assert parser is not None
        assert isinstance(parser, XlsxParser)

    def test_pptx_registered(self) -> None:
        parser = get_registry().get(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        assert parser is not None
        assert isinstance(parser, PptxParser)

    def test_generic_still_registered(self) -> None:
        # Ensure existing parsers were not displaced.
        assert get_registry().get("text/plain") is not None
        assert get_registry().get("text/markdown") is not None
